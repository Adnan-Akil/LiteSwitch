"""Contains logic for all the document conversions."""

import os
import io
import sys
import logging
import subprocess
import shutil
import platform
from typing import Optional, Callable, Dict

# Detect Linux Office Binary (LibreOffice or OpenOffice)
LINUX_OFFICE_BIN = None
if platform.system() == "Linux":
    # Prioritize libreoffice, then soffice (generic)
    LINUX_OFFICE_BIN = shutil.which("libreoffice") or shutil.which("soffice")

# Configure Logger (inherits from main app if setup, else default)
logger = logging.getLogger(__name__)

def linux_office_convert(input_path: str, output_format: str) -> Optional[str]:
    """Helper to convert using LibreOffice/OpenOffice on Linux."""
    if not LINUX_OFFICE_BIN:
         raise Exception("No Office suite found (tried 'libreoffice', 'soffice'). Please install LibreOffice.")
    
    logger.info(f"Linux Conversion: {input_path} -> {output_format} using {LINUX_OFFICE_BIN}")
    
    # Command: libreoffice --headless --convert-to pdf --outdir /path/to/dir input_file
    out_dir = os.path.dirname(input_path)
    cmd = [
        LINUX_OFFICE_BIN,
        "--headless",
        "--convert-to", output_format,
        "--outdir", out_dir,
        input_path
    ]
    
    # Run
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        logger.error(f"Office conversion failed: {result.stderr}")
        raise Exception(f"Conversion failed: {result.stderr}")
        
    # Predict output name
    base, _ = os.path.splitext(input_path)
    output_path = f"{base}.{output_format}"
    
    # Rename to adhere to LiteSwitch naming convention if needed, or just return it
    # LiteSwitch convention: name_LiteSwitch.fmt
    final_output = f"{base}_LiteSwitch.{output_format}"
    
    if os.path.exists(output_path):
        # LibreOffice output name might not match exactly what we want, rename it
        if os.path.exists(final_output):
            os.remove(final_output)
        os.rename(output_path, final_output)
        return final_output
    
    # Sometimes it might already be named correctly if we were lucky? unlikely.
    if os.path.exists(final_output):
        return final_output
        
    raise Exception("Output file not found after conversion.")

def docx_to_pdf(input_path: str) -> Optional[str]:
    """
    Converts a DOCX file to PDF.
    Windows: Uses PowerShell (COM).
    Linux: Uses LibreOffice.
    """
    if platform.system() == "Linux":
         return linux_office_convert(input_path, "pdf")
         
    """
    Converts a DOCX file to PDF using PowerShell (bypassing broken pywin32).
    Uses MS Word via COM from PowerShell.
    """
    logger.info(f"Converting DOCX to PDF (via PowerShell): {input_path}")
    
    base, _ = os.path.splitext(input_path)
    output_path = f"{base}_LiteSwitch.pdf"
    
    # PowerShell script to run Word COM
    # wdFormatPDF = 17
    ps_script = f"""
    $word = New-Object -ComObject Word.Application
    $word.Visible = $false
    try {{
        $doc = $word.Documents.Open('{input_path}')
        $doc.ExportAsFixedFormat('{output_path}', 17)
        $doc.Close($false)
    }} finally {{
        $word.Quit()
    }}
    """
    
    try:
        # Run powershell
        cmd = ["powershell", "-NoProfile", "-Command", ps_script]
        
        # Hide the PowerShell window
        startupinfo = subprocess.STARTUPINFO()
        startupinfo.dwFlags |= subprocess.STARTF_USESHOWWINDOW
        startupinfo.wShowWindow = subprocess.SW_HIDE
        
        result = subprocess.run(cmd, capture_output=True, text=True, startupinfo=startupinfo)
        
        if result.returncode != 0:
            logger.error(f"PowerShell Error: {result.stderr}")
            raise Exception(f"PowerShell conversion failed: {result.stderr}")
            
        if not os.path.exists(output_path):
             raise Exception("PDF file was not created by Word.")
             
        return output_path
        
    except Exception as e:
        logger.error(f"Error converting {input_path} to PDF: {e}")
        raise e


def docx_to_odt(input_path: str) -> Optional[str]:  # lossy
    """Converts DOCX to ODT format using pypandoc"""
    try:
        import pypandoc
        base, _ = os.path.splitext(input_path)
        output_path = f"{base}_LiteSwitch.odt"
        pypandoc.convert_file(input_path, "odt", outputfile=output_path)
        return output_path
    except Exception as e:
        logger.error(f"Error converting {input_path} to ODT: {e}")
        raise e


def docx_to_txt(input_path: str) -> Optional[str]:  # lossy
    """Converts DOCX to TXT using pypandoc"""
    try:
        import pypandoc
        base, _ = os.path.splitext(input_path)
        output_path = f"{base}_LiteSwitch.txt"
        pypandoc.convert_file(input_path, "plain", outputfile=output_path)
        return output_path
    except Exception as e:
        logger.error(f"Error converting {input_path} to TXT: {e}")
        raise e


def docx_to_md(input_path: str) -> Optional[str]:
    """Converts DOCX to Markdown using pypandoc"""
    try:
        import pypandoc
        base, _ = os.path.splitext(input_path)
        output_path = f"{base}_LiteSwitch.md"
        pypandoc.convert_file(input_path, to="markdown", outputfile=output_path)
        return output_path
    except Exception as e:
        logger.error(f"Error converting {input_path} to Markdown: {e}")
        raise e


def docx_to_latex(input_path: str) -> Optional[str]:
    """Converts DOCX to LaTeX using pypandoc"""
    try:
        import pypandoc
        base, _ = os.path.splitext(input_path)
        output_path = f"{base}_LiteSwitch.tex"
        pypandoc.convert_file(input_path, to="latex", outputfile=output_path)
        return output_path
    except Exception as e:
        logger.error(f"Error converting {input_path} to LaTeX: {e}")
        raise e


def docx_to_html(input_path: str) -> Optional[str]:  # lossy
    """Convert DOCX to HTML using pypandoc"""
    try:
        import pypandoc
        base, _ = os.path.splitext(input_path)
        output_path = f"{base}_LiteSwitch.html"
        pypandoc.convert_file(input_path, to="html", outputfile=output_path)
        return output_path
    except Exception as e:
        logger.error(f"Error converting {input_path} to HTML: {e}")
        raise e


def pdf_to_docx(input_path: str) -> Optional[str]:  # lossy
    """Convert PDF to DOCX using pdf2docx"""
    try:
        from pdf2docx import parse
        base, _ = os.path.splitext(input_path)
        output_path = f"{base}_LiteSwitch.docx"
        parse(input_path, output_path)
        return output_path
    except Exception as e:
        logger.error(f"Error converting {input_path} to DOCX: {e}")
        raise e


def pdf_to_txt(input_path: str) -> Optional[str]:
    '''Convert PDF to TXT using pdfminer.six'''
    try:
        from pdfminer.high_level import extract_text
        base, _ = os.path.splitext(input_path)
        output_path = f"{base}_LiteSwitch.txt"
        text= extract_text(input_path)
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(text)
        return output_path
    except Exception as e:
        logger.error(f"Error converting {input_path} to TXT: {e}")
        raise e


def pdf_to_png(input_path: str) -> Optional[str]: 
    '''Converts PDF to PNG using fitz'''
    try:
        import fitz
        base, _ = os.path.splitext(input_path)
        doc = fitz.open(input_path)
        last_output = ""
        for page_number, page in enumerate(doc, start=1):
            pix = page.get_pixmap(matrix=fitz.Matrix(3, 3))  # increase resolution
            output_path = f"{base}_LiteSwitch_page_{page_number}.png"
            pix.save(output_path)
            last_output = output_path # returns at least one
        doc.close()
        return last_output # returning the last one just to fit interface
    except Exception as e:
        logger.error(f"Error converting {input_path} to PNG: {e}")
        raise e


def pdf_to_html(input_path: str) -> Optional[str]:    #lossy
    '''Converts PDF to HTML using pymupdf (fitz)'''
    try:
        import fitz
        base, _ = os.path.splitext(input_path)
        output_path = f"{base}_LiteSwitch.html"
        doc = fitz.open(input_path)
        with open(output_path, 'w', encoding='utf-8') as f:
            for page in doc:
                f.write(page.get_text("html") + '\n')
        doc.close()
        return output_path
    except Exception as e:
        logger.error(f"Error converting {input_path} to HTML: {e}")
        raise e

def pdf_to_pptx(input_path: str) -> Optional[str]:
    """
    Converts PDF to PPTX using Image-Mode.
    Each page of the PDF is converted to a high-res image and placed on a slide.
    This ensures 100% visual fidelity (fonts, layout) at the cost of editability.
    """
    try:
        import fitz
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        logger.error("Missing dependencies for PPTX conversion.")
        raise
    
    try:
        logger.info(f"Converting PDF to PPTX (Image Mode): {input_path}")
        prs = Presentation()
        doc = fitz.open(input_path)

        # Set slide dimensions to match the first page of PDF (optional, but good practice)
        # For simplicity, we usually stick to default or adjust slide size.
        # Let's try to match aspect ratio of first page if possible, 
        # but changing slide size affects ALL slides in master. 
        # We'll stick to standard 16:9 or 4:3 and fit the image.
        
        for page in doc:
            # Create blank slide (layout 6 is usually blank)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Rendering high-res image (matrix=2 or 3 for better quality)
            pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0))
            img_stream = io.BytesIO(pix.tobytes("png"))
            
            # Calculate fitting
            # Powerpoint default is usually 10x7.5 inches or 13.33x7.5 (widescreen)
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            img_width_px = pix.width
            img_height_px = pix.height
            img_ratio = img_width_px / img_height_px
            slide_ratio = slide_width / slide_height
            
            # Center and Fit
            if img_ratio > slide_ratio:
                # Image is wider than slide: fit to width
                new_width = slide_width
                new_height = int(new_width / img_ratio)
                left = 0
                top = int((slide_height - new_height) / 2)
            else:
                # Image is taller/boxier: fit to height
                new_height = slide_height
                new_width = int(new_height * img_ratio)
                top = 0
                left = int((slide_width - new_width) / 2)
            
            slide.shapes.add_picture(img_stream, left, top, width=new_width, height=new_height)

        base, _ = os.path.splitext(input_path)
        output_path = f"{base}_LiteSwitch.pptx"
        prs.save(output_path)
        return output_path
    
    except Exception as e:
        logger.error(f"Error converting PDF to PPTX: {e}")
        raise e


def pdf_to_md(input_path: str) -> Optional[str]:
    '''Converts PDF to Markdown using unstructured'''
    try:
        import fitz
        from markdownify import markdownify as md
        
        base, _ = os.path.splitext(input_path)
        output_path = f"{base}_LiteSwitch.md"
        doc = fitz.open(input_path)
        markdown_content = ""

        for pageNum in range(len(doc)):
            page = doc.load_page(pageNum)
            blocks = page.get_text("blocks")

            for block in blocks:
                text = block[4].strip()
                if not text:
                    continue
                
                if len(text.splitlines()) == 1 and len(text) < 80:
                    markdown_content += f"\n## {text}\n"
                elif text.endswith('.') or text.endswith('?') or text.endswith('!'):
                    markdown_content += f"{text}\n"
                else:
                    markdown_content += f"{text}\n\n"
            
            markdown_content += "---\n\n"
            
        doc.close()
        final_markdown = md(markdown_content)
        
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(final_markdown)
        return output_path
    except Exception as e:
        logger.error(f"Error converting {input_path} to Markdown: {e}")
        raise e


def png_to_pdf(input_path: str) -> Optional[str]:
    try:
        from PIL import Image
        logger.info("opening the first image")
        first_image = Image.open(input_path).convert("RGB")
        output_path = f"{os.path.splitext(input_path)[0]}_LiteSwitch.pdf"
        first_image.save(output_path, save_all=True, dpi=(300, 300))
        return output_path
    except Exception as e:
        logger.error(f"Error converting {input_path} to PDF: {e}")
        raise e


def pptx_to_pdf(input_path: str) -> Optional[str]:
    """Converts PPTX to PDF."""
    if platform.system() == "Linux":
        return linux_office_convert(input_path, "pdf")

    """Converts PPTX to PDF using PowerShell (COM)."""
    logger.info(f"Converting PPTX to PDF: {input_path}")
    
    base, _ = os.path.splitext(input_path)
    output_path = f"{base}_LiteSwitch.pdf"
    
    # ppSaveAsPDF = 32
    ps_script = f"""
    $ppt = New-Object -ComObject PowerPoint.Application
    # PowerPoint requires at least a minimized window to work? 
    # Usually it can run hidden but sometimes needs a window.
    # We'll try minimize.
    $ppt.Visible = [Microsoft.Office.Core.MsoTriState]::msoTrue
    $ppt.WindowState = 2 # ppWindowMinimized
    
    try {{
        $pres = $ppt.Presentations.Open('{input_path}', [Microsoft.Office.Core.MsoTriState]::msoTrue, [Microsoft.Office.Core.MsoTriState]::msoFalse, [Microsoft.Office.Core.MsoTriState]::msoFalse)
        $pres.SaveAs('{output_path}', 32)
        $pres.Close()
    }} catch {{
        Write-Error $_.Exception.Message
        exit 1
    }} finally {{
        $ppt.Quit()
        [System.Runtime.Interopservices.Marshal]::ReleaseComObject($ppt) | Out-Null
    }}
    """
    
    try:
        cmd = ["powershell", "-NoProfile", "-Command", ps_script]
        
        # Hide the PowerShell window
        startupinfo = subprocess.STARTUPINFO()
        startupinfo.dwFlags |= subprocess.STARTF_USESHOWWINDOW
        startupinfo.wShowWindow = subprocess.SW_HIDE
        
        result = subprocess.run(cmd, capture_output=True, text=True, startupinfo=startupinfo)
        
        if result.returncode != 0:
             logger.error(f"PowerShell Error: {result.stderr}")
             raise Exception(f"PPTX->PDF failed: {result.stderr}")
             
        if not os.path.exists(output_path):
             raise Exception("Output PDF not found.")
             
        return output_path
    except Exception as e:
        logger.error(f"Error converting {input_path} to PDF: {e}")
        raise e

def pptx_to_png(input_path: str) -> Optional[str]:
    """Converts PPTX slides to PNG images. Returns folder path."""
    if platform.system() == "Linux":
        # Strategy: PPTX -> PDF -> PNG
        pdf_path = pptx_to_pdf(input_path)
        if pdf_path:
            # We reuse the existing pdf_to_png logic but that returns a single file path usually?
            # actually pdf_to_png in this file returns "last_output" but saves all pages.
            # We want to return a folder or similar behavior to Windows?
            # Windows implementation returns a FOLDER path.
            
            # Let's create a folder
            base, _ = os.path.splitext(input_path)
            output_dir = f"{base}_LiteSwitch_Slides"
            if not os.path.exists(output_dir):
                os.makedirs(output_dir)
            
            # Use fitz manually here to control output location matches Windows behavior
            import fitz
            doc = fitz.open(pdf_path)
            for i, page in enumerate(doc, start=1):
                pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0))
                pix.save(f"{output_dir}/Slide_{i}.png")
            doc.close()
            return output_dir

    """Converts PPTX slides to PNG images (PowerShell). Returns folder path."""
    logger.info(f"Converting PPTX to PNGs: {input_path}")
    
    base, _ = os.path.splitext(input_path)
    # Output to a folder
    output_dir = f"{base}_LiteSwitch_Slides"
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
        
    # We output to the folder. PowerPoint SaveAs PNG saves *every* slide.
    # ppSaveAsPNG = 18
    # Note: SaveAs with PNG format typically creates a folder if it's a presentation, 
    # OR we can export slide by slide. SaveAs is easiest.
    
    ps_script = f"""
    $ppt = New-Object -ComObject PowerPoint.Application
    $ppt.Visible = [Microsoft.Office.Core.MsoTriState]::msoTrue
    $ppt.WindowState = 2
    
    try {{
        $pres = $ppt.Presentations.Open('{input_path}', [Microsoft.Office.Core.MsoTriState]::msoTrue, [Microsoft.Office.Core.MsoTriState]::msoFalse, [Microsoft.Office.Core.MsoTriState]::msoFalse)
        # SaveAs with PNG(18) creates a Folder 'Filename' containing Slide1.PNG, Slide2.PNG...
        # We want to control the name, but let's see. 
        # If we target '.../Folder/Slide.png', it might save first slide?
        # Actually Export is better for slides.
        
        $i = 1
        foreach ($slide in $pres.Slides) {{
            $out = "{output_dir}\\Slide_$i.png"
            $slide.Export($out, "PNG", 1920, 1080)
            $i++
        }}
        $pres.Close()
    }} catch {{
        Write-Error $_.Exception.Message
        exit 1
    }} finally {{
        $ppt.Quit()
        [System.Runtime.Interopservices.Marshal]::ReleaseComObject($ppt) | Out-Null
    }}
    """
    
    try:
        cmd = ["powershell", "-NoProfile", "-Command", ps_script]
        
        # Hide the PowerShell window
        startupinfo = subprocess.STARTUPINFO()
        startupinfo.dwFlags |= subprocess.STARTF_USESHOWWINDOW
        startupinfo.wShowWindow = subprocess.SW_HIDE
        
        result = subprocess.run(cmd, capture_output=True, text=True, startupinfo=startupinfo)
        
        if result.returncode != 0:
             logger.error(f"PowerShell Error: {result.stderr}")
             raise Exception(f"PPTX->PNG failed: {result.stderr}")
             
        return output_dir
    except Exception as e:
        logger.error(f"Error converting {input_path} to PNG: {e}")
        raise e

def pptx_to_txt(input_path: str) -> Optional[str]:
    """Extracts text from PPTX using python-pptx."""
    if platform.system() == "Linux":
         return linux_office_convert(input_path, "txt")
         
    try:
        from pptx import Presentation
        
        base, _ = os.path.splitext(input_path)
        output_path = f"{base}_LiteSwitch.txt"
        
        prs = Presentation(input_path)
        text_runs = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
                    
        with open(output_path, "w", encoding="utf-8") as f:
            f.write("\n".join(text_runs))
            
        return output_path
    except Exception as e:
        logger.error(f"Error converting {input_path} to TXT: {e}")
        raise e

def pptx_to_docx(input_path: str) -> Optional[str]:
    """Converts PPTX text to DOCX."""
    if platform.system() == "Linux":
         return linux_office_convert(input_path, "docx")

    try:
        from pptx import Presentation
        from docx import Document
        import re

        def sanitize_xml(text):
            # Remove characters that are incompatible with XML 1.0 (control chars)
            # We keep \x09,\x0A,\x0D and normal chars. We remove \x00-\x08, \x0B, \x0C, \x0E-\x1F
            return re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F]', '', text)

        base, _ = os.path.splitext(input_path)
        output_path = f"{base}_LiteSwitch.docx"
        
        prs = Presentation(input_path)
        doc = Document()
        doc.add_heading(f"Converted from {os.path.basename(input_path)}", 0)
        
        for i, slide in enumerate(prs.slides):
            doc.add_heading(f"Slide {i+1}", level=1)
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    clean_text = sanitize_xml(shape.text)
                    if clean_text.strip():
                        doc.add_paragraph(clean_text)
        
        doc.save(output_path)
        return output_path
    except Exception as e:
        logger.error(f"Error converting {input_path} to DOCX: {e}")
        raise e


CONVERSION_MAP: Dict[str, Dict[str, Callable[[str], Optional[str]]]] = {
    "docx": {
        "pdf": docx_to_pdf,
        "odt": docx_to_odt,
        "txt": docx_to_txt,
        "md": docx_to_md,
        "tex": docx_to_latex,
    },
    "pdf": {
        "docx": pdf_to_docx,
        "png": pdf_to_png,
        "pptx": pdf_to_pptx,
        "txt": pdf_to_txt,
        "md": pdf_to_md,
    },
    "pptx": {
        "pdf": pptx_to_pdf,
        "png": pptx_to_png,
        "txt": pptx_to_txt,
        "docx": pptx_to_docx,
    },
    "ppt": {
        "pdf": pptx_to_pdf,
        "png": pptx_to_png,
        "txt": pptx_to_txt,
        "docx": pptx_to_docx,
    },
    "png":{
        "pdf": png_to_pdf
    }
}
